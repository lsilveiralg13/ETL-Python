import os
from pptx import Presentation
from pptx.util import Cm

# =========================================================
# 📂 Caminhos
# =========================================================
# Pasta onde estão as imagens exportadas
pasta_imagens = r"C:\Users\lucasbarros\OneDrive - CTC FRANCHISING S A\Área de Trabalho\Cards_Exportados"

# Caminho de saída do PowerPoint final
ppt_saida = r"C:\Users\lucasbarros\OneDrive - CTC FRANCHISING S A\Área de Trabalho\Cards_Novas_Lojas_Final.pptx"

# =========================================================
# 🧱 Parâmetros de layout
# =========================================================
largura_card = Cm(10.8)
altura_card = Cm(9.3)
espaco_horizontal = Cm(0.5)
espaco_vertical = Cm(0.5)
margem_esquerda = Cm(0.4)
margem_superior = Cm(0.4)

# =========================================================
# 🖼️ Carrega as imagens
# =========================================================
if not os.path.exists(pasta_imagens):
    raise FileNotFoundError(f"A pasta '{pasta_imagens}' não foi encontrada.")

# Filtra somente imagens válidas
imagens = sorted([
    os.path.join(pasta_imagens, f)
    for f in os.listdir(pasta_imagens)
    if f.lower().endswith((".png", ".jpg", ".jpeg"))
])

if not imagens:
    raise ValueError("Nenhuma imagem foi encontrada na pasta especificada.")

print(f"🔍 {len(imagens)} cards encontrados para montagem.")

# =========================================================
# 🧩 Criação da apresentação final
# =========================================================
prs = Presentation()

for i in range(0, len(imagens), 6):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subset = imagens[i:i+6]

    for idx, imagem in enumerate(subset):
        linha = idx // 3  # 0 ou 1
        coluna = idx % 3  # 0, 1, 2

        left = margem_esquerda + coluna * (largura_card + espaco_horizontal)
        top = margem_superior + linha * (altura_card + espaco_vertical)

        slide.shapes.add_picture(imagem, left, top, width=largura_card, height=altura_card)

# =========================================================
# 💾 Salvamento do arquivo
# =========================================================
prs.save(ppt_saida)
print(f"\n🎉 Apresentação final gerada com sucesso!")
print(f"📂 Caminho do arquivo: {ppt_saida}")
